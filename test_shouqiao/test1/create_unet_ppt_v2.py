from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ================================================================
# 复用参考PPT的风格：16:9 / #3667E8 主色 / Century Gothic 数字 /
# 微软雅黑中文 / 卡片式布局 / 干净留白
# ================================================================

prs = Presentation()
prs.slide_width = 12192000   # 16:9
prs.slide_height = 6858000

# ===== 调色板 =====
BLUE    = RGBColor(0x36, 0x67, 0xE8)   # 主色
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深色标题
GRAY    = RGBColor(0x66, 0x66, 0x66)   # 正文灰
LGRAY   = RGBColor(0xF5, 0xF5, 0xF8)   # 背景浅灰
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF0, 0xF4, 0xFC)   # 卡片浅蓝底
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
RED     = RGBColor(0xC6, 0x28, 0x28)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
TEAL    = RGBColor(0x00, 0x89, 0x7B)
PURPLE  = RGBColor(0x7B, 0x2D, 0x8B)

# ===== 工具函数 =====
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_c, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_c
    s.line.fill.background()
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    return s

def add_rounded(slide, l, t, w, h, fill_c, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_c
    s.line.fill.background()
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    # set corner radius
    sp = s._element
    spPr = sp.find(qn('a:spPr')) if sp.find(qn('a:spPr')) is not None else sp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is not None:
        prstGeom = spPr.find(qn('a:prstGeom')) if spPr.find(qn('a:prstGeom')) is not None else spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                from lxml import etree
                avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
                gd.set('name', 'adj')
                gd.set('fmla', 'val 5000')
    return s

def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def set_p(tf, text, sz=18, clr=DARK, bold=False, align=PP_ALIGN.LEFT, sa=Pt(6), sb=Pt(0), font_name=None):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    p.space_after = sa
    p.space_before = sb
    if font_name:
        p.font.name = font_name
        # Set East Asian font
        for run in p.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                from lxml import etree
                rPr = etree.SubElement(run._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                from lxml import etree
                ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
            ea.set('typeface', font_name)
    return p

def add_p(tf, text, sz=18, clr=DARK, bold=False, align=PP_ALIGN.LEFT, sa=Pt(4), sb=Pt(0), font_name=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    p.space_after = sa
    p.space_before = sb
    if font_name:
        p.font.name = font_name
        for run in p.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                from lxml import etree
                rPr = etree.SubElement(run._r, '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                from lxml import etree
                ea = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
            ea.set('typeface', font_name)
    return p

def make_section_slide(prs, number, title, bg_color=WHITE):
    """创建章节分隔页：左侧大数字 + 右侧标题"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, bg_color)
    # 顶部蓝色装饰条
    add_rect(s, 0, 0, prs.slide_width, Emu(55000), BLUE)
    # 大数字
    tb_num = add_tb(s, Emu(2903855), Emu(1968500), Emu(6384290), Emu(1007745))
    set_p(tb_num.text_frame, number, sz=28, clr=BLUE, bold=True, font_name='Century Gothic')
    # 标题
    tb_title = add_tb(s, Emu(2906395), Emu(3159125), Emu(6384290), Emu(1051560))
    set_p(tb_title.text_frame, title, sz=28, clr=DARK, bold=True, font_name='微软雅黑')
    return s

def make_content_header(slide, title):
    """内容页标题栏：顶部蓝色条 + 左侧标题"""
    add_rect(slide, 0, 0, prs.slide_width, Emu(55000), BLUE)
    tb = add_tb(slide, Emu(92075), Emu(179705), Emu(10516870), Emu(600000))
    set_p(tb.text_frame, title, sz=24, clr=DARK, bold=True, font_name='微软雅黑')
    # 底部装饰线
    add_rect(slide, Emu(92075), Emu(780000), Emu(2000000), Emu(35000), BLUE)

# ================================================================
# SLIDE 1 — 封面
# ================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
# 顶部大色块
add_rect(s1, 0, 0, prs.slide_width, Emu(3200000), RGBColor(0xF0, 0xF4, 0xFC))
# 蓝色装饰条
add_rect(s1, 0, Emu(3200000), prs.slide_width, Emu(55000), BLUE)
# 主标题
tb = add_tb(s1, Emu(1058737), Emu(1200000), Emu(10000000), Emu(800000))
set_p(tb.text_frame, "从 U-Net 到 nnU-Net", sz=40, clr=DARK, bold=True, font_name='微软雅黑')
# 副标题
tb2 = add_tb(s1, Emu(1058737), Emu(2050000), Emu(10000000), Emu(500000))
set_p(tb2.text_frame, "医学图像分割的标准化流水线", sz=22, clr=GRAY, font_name='微软雅黑')
# 底部信息
tb3 = add_tb(s1, Emu(1058737), Emu(3600000), Emu(10000000), Emu(400000))
set_p(tb3.text_frame, "组会学术汇报  |  2025年4月17日", sz=16, clr=GRAY, font_name='微软雅黑')
# 底部渐变装饰
add_rect(s1, 0, Emu(6500000), prs.slide_width, Emu(358000), BLUE)

# ================================================================
# SLIDE 2 — 目录 (参考PPT的网格风格)
# ================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, WHITE)
add_rect(s2, 0, 0, prs.slide_width, Emu(55000), BLUE)
tb_title = add_tb(s2, Emu(92075), Emu(179705), Emu(5000000), Emu(600000))
set_p(tb_title.text_frame, "CONTENT", sz=28, clr=BLUE, bold=True, font_name='Century Gothic')
add_rect(s2, Emu(92075), Emu(780000), Emu(2000000), Emu(35000), BLUE)

toc_items = [
    ("01", "U-Net 回顾与缺陷"),
    ("02", "nnU-Net 提出动机"),
    ("03", "网络架构微操"),
    ("04", "动态自适应机制"),
    ("05", "训练策略与后处理"),
    ("06", "级联架构与实验总结"),
]
positions = [
    (Emu(1655445), Emu(1454150)),
    (Emu(4552315), Emu(1454150)),
    (Emu(7449820), Emu(1454150)),
    (Emu(1655445), Emu(3966210)),
    (Emu(4552315), Emu(3966210)),
    (Emu(7449820), Emu(3966210)),
]

for (num, title), (l, t) in zip(toc_items, positions):
    card = add_rounded(s2, l, t, Emu(2897505), Emu(2512060), CARD_BG, BLUE, Pt(1))
    # 数字
    tb_num = add_tb(s2, l + Emu(200000), t + Emu(800000), Emu(1300000), Emu(500000))
    set_p(tb_num.text_frame, num, sz=36, clr=BLUE, bold=True, font_name='Century Gothic')
    # 标题
    tb_t = add_tb(s2, l + Emu(200000), t + Emu(1500000), Emu(2500000), Emu(600000))
    set_p(tb_t.text_frame, title, sz=18, clr=DARK, bold=True, font_name='微软雅黑')

# ================================================================
# SLIDE 3 — 章节：U-Net回顾
# ================================================================
s3 = make_section_slide(prs, "01", "U-Net 回顾与缺陷")

# ================================================================
# SLIDE 4 — U-Net缺陷 (内容页)
# ================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, WHITE)
make_content_header(s4, "U-Net 回顾与缺陷")

# 左侧优点
box_l = add_rounded(s4, Emu(300000), Emu(1100000), Emu(5600000), Emu(5200000), CARD_BG, BLUE, Pt(1.5))
tb_l = add_tb(s4, Emu(500000), Emu(1250000), Emu(5200000), Emu(5000000))
tf = tb_l.text_frame
set_p(tf, "U-Net 经典贡献", sz=22, clr=BLUE, bold=True, sa=Pt(12), font_name='微软雅黑')
for line in [
    "• U 型对称结构（编码器-解码器）",
    "• 跳跃连接融合底层细节与高层语义",
    "• 数据增强驱动，少量标注即可训练",
    "• 医学图像分割的奠基之作",
]:
    add_p(tf, line, sz=16, clr=DARK, sa=Pt(6), font_name='微软雅黑')

# 右侧缺陷
box_r = add_rounded(s4, Emu(6200000), Emu(1100000), Emu(5700000), Emu(5200000), RGBColor(0xFC, 0xE4, 0xE4), RED, Pt(1.5))
tb_r = add_tb(s4, Emu(6400000), Emu(1250000), Emu(5300000), Emu(5000000))
tf2 = tb_r.text_frame
set_p(tf2, "原始 U-Net 的局限", sz=22, clr=RED, bold=True, sa=Pt(12), font_name='微软雅黑')
for line in [
    "❶ 空间维度灾难",
    "    2D 看不到 Z 轴；3D 撑爆显存",
    "❷ 数据特性鸿沟",
    "    CT(HU 值) vs MRI(相对亮度) 一刀切",
    "❸ 静态模型，无法自适应",
    "❹ 效果不好就魔改网络",
    "    忽略预处理和超参数调优",
]:
    add_p(tf2, line, sz=16, clr=DARK, sa=Pt(4), font_name='微软雅黑')

# ================================================================
# SLIDE 5 — 章节：nnU-Net提出动机
# ================================================================
s5 = make_section_slide(prs, "02", "nnU-Net 提出动机")

# ================================================================
# SLIDE 6 — nnU-Net 核心思想 (内容页)
# ================================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, WHITE)
make_content_header(s6, "nnU-Net：no-new-Net 理念")

# 核心引语
box_q = add_rounded(s6, Emu(300000), Emu(1100000), Emu(11400000), Emu(1400000), CARD_BG, BLUE, Pt(1.5))
tb_q = add_tb(s6, Emu(500000), Emu(1200000), Emu(11000000), Emu(1200000))
tf = tb_q.text_frame
set_p(tf, "核心思想", sz=20, clr=BLUE, bold=True, sa=Pt(8), font_name='微软雅黑')
add_p(tf, "一个没有任何花哨模块的最基础 U-Net，只要把数据预处理、网络拓扑生成、训练策略和推理后处理", sz=16, clr=DARK, font_name='微软雅黑')
add_p(tf, "做到极致，就能在 MSD 十项全能挑战赛中打败所有当时最先进的魔改网络。", sz=16, clr=DARK, font_name='微软雅黑')
add_p(tf, "→ 把深度学习变成标准化的工业流水线", sz=17, clr=BLUE, bold=True, font_name='微软雅黑')

# 四张卡片
cards = [
    ("数据预处理", "自动分析分辨率、大小、\n强度分布，针对性归一化", BLUE),
    ("网络拓扑生成", "根据显存与图像大小\n自动推导 Patch 与深度", TEAL),
    ("训练策略", "Dice+CE 组合损失\n动态调整超参数", ORANGE),
    ("推理后处理", "重叠滑窗 + 高斯加权\n+ 连通域清洗", PURPLE),
]
for i, (title, desc, color) in enumerate(cards):
    l = Emu(300000 + i * 2920000)
    t = Emu(2800000)
    card = add_rounded(s6, l, t, Emu(2720000), Emu(3400000), WHITE, color, Pt(1.5))
    # 编号圆圈
    circ = add_rounded(s6, l + Emu(100000), t + Emu(150000), Emu(400000), Emu(400000), color)
    tb_c = add_tb(s6, l + Emu(100000), t + Emu(150000), Emu(400000), Emu(400000))
    set_p(tb_c.text_frame, str(i+1), sz=18, clr=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Century Gothic')
    # 标题
    tb_t = add_tb(s6, l + Emu(600000), t + Emu(180000), Emu(2000000), Emu(400000))
    set_p(tb_t.text_frame, title, sz=20, clr=color, bold=True, font_name='微软雅黑')
    # 描述
    tb_d = add_tb(s6, l + Emu(150000), t + Emu(800000), Emu(2400000), Emu(2400000))
    set_p(tb_d.text_frame, desc, sz=15, clr=GRAY, font_name='微软雅黑')

# ================================================================
# SLIDE 7 — 章节：架构微操
# ================================================================
s7 = make_section_slide(prs, "03", "网络架构微操与动态自适应")

# ================================================================
# SLIDE 8 — 架构改进内容页
# ================================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, WHITE)
make_content_header(s8, "网络架构的『微操』与『变形』")

# 三个卡片
arch_items = [
    ("ReLU → Leaky ReLU", "原始 U-Net 初始化精妙，Dying ReLU\n不严重。但 nnU-Net 通杀所有医学图像，\n面对低对比度数据时易出现大量\n神经元进入负区间，需 Leaky ReLU\n让负区间仍有梯度流动。", BLUE),
    ("BN → Instance Norm", "3D 训练时 Batch Size 极小（通常 2），\nBN 统计量崩溃。Instance Norm 对\n单个三维块归一化，不受 batch size\n影响，适合 3D 医疗图像训练。", TEAL),
    ("三种自适应形态", "▸ 2D U-Net：适合薄层切片\n▸ 3D U-Net：适合各向同性体数据\n▸ 3D U-Net Cascade：适合大图像\n+ 高精度需求（低分辨率全局 +\n高分辨率局部）", ORANGE),
]
for i, (title, desc, color) in enumerate(arch_items):
    l = Emu(300000 + i * 3920000)
    t = Emu(1300000)
    card = add_rounded(s8, l, t, Emu(3720000), Emu(5000000), CARD_BG, color, Pt(1.5))
    tb_t = add_tb(s8, l + Emu(150000), t + Emu(200000), Emu(3400000), Emu(500000))
    set_p(tb_t.text_frame, f"0{i+1}  {title}", sz=20, clr=color, bold=True, font_name='微软雅黑')
    tb_d = add_tb(s8, l + Emu(150000), t + Emu(900000), Emu(3400000), Emu(3800000))
    set_p(tb_d.text_frame, desc, sz=15, clr=DARK, font_name='微软雅黑')

# ================================================================
# SLIDE 9 — 训练策略与后处理
# ================================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, WHITE)
make_content_header(s9, "训练策略与后处理")

# 左侧：损失函数
box_l = add_rounded(s9, Emu(300000), Emu(1100000), Emu(5600000), Emu(5200000), CARD_BG, BLUE, Pt(1.5))
tb_l = add_tb(s9, Emu(500000), Emu(1250000), Emu(5200000), Emu(5000000))
tf = tb_l.text_frame
set_p(tf, "组合损失函数：Dice + Cross Entropy", sz=20, clr=BLUE, bold=True, sa=Pt(10), font_name='微软雅黑')
add_p(tf, "Dice Loss", sz=18, clr=GREEN, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf, "• 衡量预测区域与真实区域的重合度", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 整体性全局监督，天生抗不平衡", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 适合小目标分割（如血管）", sz=15, clr=DARK, sa=Pt(10), font_name='微软雅黑')
add_p(tf, "Cross Entropy Loss", sz=18, clr=RED, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf, "• 逐像素分类监督", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 训练稳定，收敛平滑", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 但偏向背景，前景过小时被主导", sz=15, clr=DARK, sa=Pt(10), font_name='微软雅黑')
add_p(tf, "二者互补 → 分割精度显著提升", sz=16, clr=BLUE, bold=True, font_name='微软雅黑')

# 右侧：后处理
box_r = add_rounded(s9, Emu(6200000), Emu(1100000), Emu(5700000), Emu(5200000), RGBColor(0xFD, 0xF0, 0xE8), ORANGE, Pt(1.5))
tb_r = add_tb(s9, Emu(6400000), Emu(1250000), Emu(5300000), Emu(5000000))
tf2 = tb_r.text_frame
set_p(tf2, "推理后处理三件套", sz=20, clr=ORANGE, bold=True, sa=Pt(10), font_name='微软雅黑')
add_p(tf2, "① 50% 重叠滑窗", sz=17, clr=ORANGE, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf2, "前一个块的边缘 = 后一个块的中心", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf2, "中心视野好，预测精度极高", sz=15, clr=DARK, sa=Pt(8), font_name='微软雅黑')
add_p(tf2, "② 高斯加权融合", sz=17, clr=TEAL, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf2, "中心高权重(~0.99)、边缘低权重(~0.05)", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf2, "用自信区域覆盖边缘误差", sz=15, clr=DARK, sa=Pt(8), font_name='微软雅黑')
add_p(tf2, "③ 连通域后处理", sz=17, clr=PURPLE, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf2, "保留最大连通域，消除孤立噪点块", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')

# ================================================================
# SLIDE 10 — Cascade 架构 + 总结
# ================================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, WHITE)
make_content_header(s10, "级联架构 (Cascade) 与实验总结")

# 左侧：Cascade
box_c = add_rounded(s10, Emu(300000), Emu(1100000), Emu(5600000), Emu(5200000), CARD_BG, BLUE, Pt(1.5))
tb_c = add_tb(s10, Emu(500000), Emu(1250000), Emu(5200000), Emu(5000000))
tf = tb_c.text_frame
set_p(tf, "3D U-Net Cascade", sz=20, clr=BLUE, bold=True, sa=Pt(10), font_name='微软雅黑')
add_p(tf, "Stage 1：低分辨率全局定位", sz=17, clr=TEAL, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf, "• 原始高清图降采样 → 看全局", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 输出『粗略导航地图』", sz=15, clr=DARK, sa=Pt(8), font_name='微软雅黑')
add_p(tf, "Stage 2：高分辨率局部精细分割", sz=17, clr=ORANGE, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf, "• 高清 Patch + Stage1 导航图拼接", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf, "• 双重输入 → 精密局部分割", sz=15, clr=DARK, sa=Pt(8), font_name='微软雅黑')
add_p(tf, "类比：远处看全景 → 近处放大镜+ GPS", sz=16, clr=DARK, bold=True, font_name='微软雅黑')

# 右侧：总结
box_s = add_rounded(s10, Emu(6200000), Emu(1100000), Emu(5700000), Emu(5200000), RGBColor(0xE8, 0xF5, 0xE9), GREEN, Pt(1.5))
tb_s = add_tb(s10, Emu(6400000), Emu(1250000), Emu(5300000), Emu(5000000))
tf2 = tb_s.text_frame
set_p(tf2, "MSD 结果与总结启示", sz=20, clr=GREEN, bold=True, sa=Pt(10), font_name='微软雅黑')
add_p(tf2, "MSD 十项全能挑战赛", sz=17, clr=DARK, bold=True, sa=Pt(6), font_name='微软雅黑')
add_p(tf2, "nnU-Net 以全自动最基础 U-Net", sz=15, clr=DARK, sa=Pt(2), font_name='微软雅黑')
add_p(tf2, "击败所有当时最先进魔改网络", sz=15, clr=DARK, sa=Pt(10), font_name='微软雅黑')
add_p(tf2, "启示", sz=17, clr=GREEN, bold=True, sa=Pt(6), font_name='微软雅黑')
for line in [
    "❶ 效果不好 ≠ 网络不够先进",
    "❷ 先检查数据重采样 / 归一化 / Patch",
    "❸ 底层工程做到极致",
    "❹ 建立坚如磐石的 Baseline",
    "❺ 工程 > 魔改",
]:
    add_p(tf2, line, sz=15, clr=DARK, sa=Pt(3), font_name='微软雅黑')

# ================================================================
# SAVE
# ================================================================
out = "d:/PythonProject/hello_agent/fastapi_/从U-Net到nnU-Net_组会汇报.pptx"
prs.save(out)
print(f"PPT saved: {out}")
