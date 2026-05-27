from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 调色板 =====
DARK_NAVY   = RGBColor(0x0D, 0x23, 0x3F)
MED_BLUE    = RGBColor(0x1B, 0x4F, 0x8A)
LIGHT_BLUE  = RGBColor(0xD6, 0xE6, 0xF5)
ACCENT_TEAL = RGBColor(0x00, 0x89, 0x7B)
ACCENT_GRN  = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED  = RGBColor(0xC6, 0x28, 0x28)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
BG_CREAM    = RGBColor(0xFA, 0xFA, 0xFA)

# ===== 工具函数 =====
def set_bg(slide, color=BG_CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_c, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_c
    s.line.fill.background()
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    return s

def add_rounded(slide, l, t, w, h, fill_c, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_c
    s.line.fill.background()
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    return s

def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def set_p(tf, text, sz=18, clr=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, sa=Pt(6), sb=Pt(0)):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    p.space_after = sa
    p.space_before = sb
    return p

def add_p(tf, text, sz=18, clr=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, sa=Pt(4), sb=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    p.space_after = sa
    p.space_before = sb
    return p

def page_header(slide, title, subtitle=None):
    """统一页眉：顶部色条 + 标题 + 装饰线"""
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.07), MED_BLUE)
    tb = add_tb(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.6))
    set_p(tb.text_frame, title, 28, DARK_NAVY, bold=True)
    add_rect(slide, Inches(0.8), Inches(0.95), Inches(2.2), Inches(0.04), ACCENT_TEAL)
    if subtitle:
        tb2 = add_tb(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.35))
        set_p(tb2.text_frame, subtitle, 15, MED_GRAY)


# ================================================================
# SLIDE 1 — 封面
# ================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, WHITE)
add_rect(s1, Inches(0), Inches(0), Inches(13.333), Inches(0.1), DARK_NAVY)

# 左侧装饰竖条
add_rect(s1, Inches(1.0), Inches(2.0), Inches(0.07), Inches(3.0), ACCENT_TEAL)

tb = add_tb(s1, Inches(1.4), Inches(2.0), Inches(10), Inches(1.2))
set_p(tb.text_frame, "从 U-Net 到 nnU-Net", 42, DARK_NAVY, bold=True)

tb2 = add_tb(s1, Inches(1.4), Inches(3.2), Inches(10), Inches(0.7))
set_p(tb2.text_frame, "医学图像分割的标准化流水线——从人工调参到自动适配", 22, MED_GRAY)

tb3 = add_tb(s1, Inches(1.4), Inches(4.2), Inches(10), Inches(0.6))
set_p(tb3.text_frame, "组会学术汇报  |  2025年4月17日", 16, MED_GRAY)

# 底部装饰
add_rect(s1, Inches(0), Inches(7.4), Inches(13.333), Inches(0.1), DARK_NAVY)

# ================================================================
# SLIDE 2 — 目录
# ================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, WHITE)
page_header(s2, "汇报提纲")

items = [
    ("01", "U-Net 回顾：经典架构与固有缺陷"),
    ("02", "nnU-Net 的提出动机：no-new-Net 理念"),
    ("03", "网络架构微操：Leaky ReLU / IN / 三种形态"),
    ("04", "动态自适应机制：让网络自己算参数"),
    ("05", "训练策略：组合损失函数设计"),
    ("06", "推理与后处理：无缝拼接的艺术"),
    ("07", "级联架构 (Cascade) 详解"),
    ("08", "实验结果与总结启示"),
]

for i, (num, title) in enumerate(items):
    row = i // 2
    col = i % 2
    l = Inches(0.8 + col * 6.2)
    t = Inches(1.6 + row * 1.3)
    box = add_rounded(s2, l, t, Inches(5.5), Inches(1.0), RGBColor(0xF0, 0xF4, 0xFA), MED_BLUE, Pt(1))
    # 编号
    tb_num = add_tb(s2, Inches(0.9 + col * 6.2), Inches(1.7 + row * 1.3), Inches(0.6), Inches(0.7))
    set_p(tb_num.text_frame, num, 24, MED_BLUE, bold=True, align=PP_ALIGN.CENTER)
    # 标题
    tb_t = add_tb(s2, Inches(1.6 + col * 6.2), Inches(1.75 + row * 1.3), Inches(4.5), Inches(0.7))
    set_p(tb_t.text_frame, title, 17, DARK_GRAY)

# ================================================================
# SLIDE 3 — U-Net 回顾与缺陷
# ================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, WHITE)
page_header(s3, "U-Net 回顾：经典架构与固有缺陷", "医学图像分割的奠基之作，但在实际应用中暴露致命局限")

# 左：U-Net 优点
box_l = add_rounded(s3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GRN, Pt(1))
tb_l = add_tb(s3, Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.8))
tf = tb_l.text_frame
set_p(tf, "U-Net 经典贡献", 22, ACCENT_GRN, bold=True, sa=Pt(12))
for line in [
    "• U 型对称结构：编码器-解码器架构",
    "• 跳跃连接 (Skip Connection)：",
    "  融合底层细节与高层语义",
    "• 数据增强驱动：少量标注即可训练",
    "• 医学图像分割的标杆与基石",
    "",
    "  优势",
    "  ✓ 在小样本医学数据上表现出色",
    "  ✓ 结构简洁，易于理解与实现",
    "  ✓ 后续所有分割网络的起点",
]:
    add_p(tf, line, 15, DARK_GRAY, sa=Pt(4))

# 右：缺陷
box_r = add_rounded(s3, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2), RGBColor(0xFC, 0xE4, 0xE4), ACCENT_RED, Pt(1))
tb_r = add_tb(s3, Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8))
tf2 = tb_r.text_frame
set_p(tf2, "原始 U-Net 的致命缺陷", 22, ACCENT_RED, bold=True, sa=Pt(12))
for line in [
    "❶ 空间维度灾难",
    "   2D 看不到 Z 轴；3D 直接撑爆显存",
    "",
    "❷ 数据特性鸿沟",
    "   CT 有绝对 HU 值 vs MRI 仅有相对亮度",
    "   各向同性 vs 各向异性，一刀切地处理",
    "",
    "❸ 静态模型，无法自适应",
    "   对不同分辨率/模态的数据用同一套参数",
    "",
    "❹ 科研界的『弯路』",
    "   效果不好就魔改网络，",
    "   却忽略了预处理和超参数调优",
]:
    add_p(tf2, line, 15, DARK_GRAY, sa=Pt(4))

# ================================================================
# SLIDE 4 — nnU-Net 提出动机
# ================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, WHITE)
page_header(s4, "nnU-Net 的提出动机", "no-new-Net —— 没有新网络，只有新流程")

# 核心思想大框
box_main = add_rounded(s4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2), RGBColor(0xE8, 0xF0, 0xF8), MED_BLUE, Pt(1.5))
tb_main = add_tb(s4, Inches(1.2), Inches(1.6), Inches(11.0), Inches(2.0))
tf = tb_main.text_frame
set_p(tf, "核心思想", 24, DARK_NAVY, bold=True, sa=Pt(10))
add_p(tf, "一个没有任何花哨模块的、最基础的 U-Net，只要把数据预处理、网络拓扑生成、训练策略和推理后处理", 17, DARK_GRAY, sa=Pt(2))
add_p(tf, "这四个周边环节做到极致，就能在十项全能（MSD）挑战赛中打败所有当时最先进的魔改网络。", 17, DARK_GRAY, sa=Pt(2))
add_p(tf, "它把深度学习变成了一条标准化的工业流水线。", 17, MED_BLUE, bold=True, sa=Pt(8))

# 四个环节
boxes_data = [
    ("数据预处理", "自动分析分辨率、图像大小、\n强度分布，针对性归一化", RGBColor(0x2C, 0x5F, 0x8A)),
    ("网络拓扑生成", "根据显存限制与中位数图像\n大小自动推推导 Patch 与深度", RGBColor(0x00, 0x89, 0x7B)),
    ("训练策略", "组合损失函数(Dice+CE)、\n动态调整训练超参数", RGBColor(0xE6, 0x7E, 0x22)),
    ("推理后处理", "重叠滑窗 + 高斯加权融合\n+ 连通域清洗", RGBColor(0x7B, 0x2D, 0x8B)),
]
for i, (title, desc, color) in enumerate(boxes_data):
    l = Inches(0.8 + i * 3.1)
    t = Inches(4.2)
    box = add_rounded(s4, l, t, Inches(2.8), Inches(2.8), RGBColor(0xFA, 0xFA, 0xFA), color, Pt(1.5))
    # 序号圆点
    dot = add_rounded(s4, l + Inches(0.1), t + Inches(0.15), Inches(0.35), Inches(0.35), color)
    tb_dot = add_tb(s4, l + Inches(0.1), t + Inches(0.15), Inches(0.35), Inches(0.35))
    set_p(tb_dot.text_frame, str(i+1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 标题
    tb_t = add_tb(s4, l + Inches(0.55), t + Inches(0.15), Inches(2.1), Inches(0.4))
    set_p(tb_t.text_frame, title, 17, color, bold=True)
    # 描述
    tb_d = add_tb(s4, l + Inches(0.15), t + Inches(0.65), Inches(2.5), Inches(2.0))
    set_p(tb_d.text_frame, desc, 14, DARK_GRAY, sa=Pt(4))

# ================================================================
# SLIDE 5 — 架构微操
# ================================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, WHITE)
page_header(s5, "网络架构的『微操』与『变形』", "不改网络主体结构，只优化细节")

# 左：Leaky ReLU
box1 = add_rounded(s5, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.3), RGBColor(0xE8, 0xF0, 0xF8), MED_BLUE, Pt(1))
tb1 = add_tb(s5, Inches(1.0), Inches(1.7), Inches(3.3), Inches(5.0))
tf = tb1.text_frame
set_p(tf, "① ReLU → Leaky ReLU", 18, DARK_NAVY, bold=True, sa=Pt(10))
for line in [
    "原始 U-Net 用 ReLU，",
    "初始化极其精妙，训练初期",
    "神经元均处『活跃』态，",
    "没有严重的 Dying ReLU。",
    "",
    "但 nnU-Net 目标是『通杀』",
    "所有医学图像，面对复杂、",
    "低对比度数据时，大量神经",
    "元可能进入负区间。",
    "",
    "→ Leaky ReLU 让负区间",
    "  仍有微小梯度流动",
]:
    add_p(tf, line, 14, DARK_GRAY, sa=Pt(3))

# 中：BN → IN
box2 = add_rounded(s5, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.3), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GRN, Pt(1))
tb2 = add_tb(s5, Inches(5.0), Inches(1.7), Inches(3.3), Inches(5.0))
tf2 = tb2.text_frame
set_p(tf2, "② BN → Instance Norm", 18, DARK_NAVY, bold=True, sa=Pt(10))
for line in [
    "Batch Normalization (BN)",
    "在 3D 训练时 Batch Size",
    "极小（通常仅 2），",
    "BN 统计量崩溃。",
    "",
    "→ Instance Normalization (IN)",
    "  对单个三维块进行归一化",
    "",
    "IN 优势：",
    "✓ 不受 batch size 影响",
    "✓ 适合 3D 医疗图像",
    "✓ 稳定训练过程",
]:
    add_p(tf2, line, 14, DARK_GRAY, sa=Pt(3))

# 右：三种形态
box3 = add_rounded(s5, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.3), RGBColor(0xFD, 0xF0, 0xE8), RGBColor(0xE6, 0x7E, 0x22), Pt(1))
tb3 = add_tb(s5, Inches(9.0), Inches(1.7), Inches(3.3), Inches(5.0))
tf3 = tb3.text_frame
set_p(tf3, "③ 三种自适应形态", 18, DARK_NAVY, bold=True, sa=Pt(10))
for line in [
    "nnU-Net 自动选择最优方案：",
    "",
    "▸ 2D U-Net",
    "  适合薄层切片数据",
    "",
    "▸ 3D U-Net",
    "  适合各向同性体数据",
    "",
    "▸ 3D U-Net Cascade",
    "  适合大图像 + 高精度需求",
    "  (低分辨率全局 + 高分辨率局部)",
]:
    add_p(tf3, line, 14, DARK_GRAY, sa=Pt(3))

# ================================================================
# SLIDE 6 — 动态自适应
# ================================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, WHITE)
page_header(s6, "动态自适应机制", "nnU-Net 不直接训练，而是先『读懂』你的数据集")

# 流程步骤
steps = [
    ("分析数据集", ["分辨率", "图像大小", "强度分布"], ACCENT_TEAL),
    ("自动归一化", ["CT：HU 值统计→截断+Z-score", "MRI：单例内部归一化"], MED_BLUE),
    ("反推网络拓扑", ["固定显存上限（如 12GB）", "计算最大 Patch 大小", "确定池化层数（→ 8×8×8）"], RGBColor(0xE6, 0x7E, 0x22)),
    ("动态配置", ["自动设定 Batch Size", "自动选择 2D/3D/Cascade"], RGBColor(0x7B, 0x2D, 0x8B)),
]
for i, (title, descs, color) in enumerate(steps):
    l = Inches(0.8 + i * 3.1)
    t = Inches(1.5)
    # 箭头连接
    if i < 3:
        arrow = add_tb(s6, l + Inches(2.65), t + Inches(0.5), Inches(0.7), Inches(0.3))
        set_p(arrow.text_frame, "→", 24, MED_GRAY, align=PP_ALIGN.CENTER)

    box = add_rounded(s6, l, t, Inches(2.85), Inches(2.0), color)
    tb_t = add_tb(s6, l + Inches(0.1), t + Inches(0.1), Inches(2.65), Inches(0.4))
    set_p(tb_t.text_frame, f"Step {i+1}", 14, WHITE, bold=True)
    tb_t2 = add_tb(s6, l + Inches(0.1), t + Inches(0.45), Inches(2.65), Inches(0.4))
    set_p(tb_t2.text_frame, title, 20, WHITE, bold=True, sa=Pt(2))
    tb_d = add_tb(s6, l + Inches(0.1), t + Inches(1.0), Inches(2.65), Inches(1.0))
    tf = tb_d.text_frame
    for d in descs:
        add_p(tf, "• " + d, 13, RGBColor(0xFF, 0xFF, 0xEE), sa=Pt(2))

# 底部：具体示例
box_ex = add_rounded(s6, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.8), RGBColor(0xF5, 0xF5, 0xF5), MED_BLUE, Pt(1))
tb_ex = add_tb(s6, Inches(1.0), Inches(4.0), Inches(11.3), Inches(2.5))
tf = tb_ex.text_frame
set_p(tf, "实际案例：nnU-Net 根据数据集自动生成的参数", 18, DARK_NAVY, bold=True, sa=Pt(10))

header_fmt = "{:<20} {:<18} {:<18} {:<18}"
row_fmt = "{:<20} {:<18} {:<18} {:<18}"

lines = [
    f"{'数据集':<14} {'中位数形状':<18} {'Patch 大小':<18} {'各轴池化次数':<18}",
    f"{'BrainTumour':<14} {'138 × 169 × 138':<18} {'128 × 128 × 128':<18} {'5, 5, 5':<18}",
    f"{'Liver':<14} {'512 × 512 × 367':<18} {'128 × 128 × 128':<18} {'5, 5, 5':<18}",
    f"{'Lung':<14} {'512 × 512 × 271':<18} {'80 × 80 × 192':<18} {'5, 5, 6':<18}",
    f"{'Pancreas':<14} {'512 × 512 × 153':<18} {'40 × 224 × 224':<18} {'4, 5, 5':<18}",
]
for line in lines:
    add_p(tf, line, 13, DARK_GRAY, sa=Pt(3))

add_p(tf, "\n核心结论：没有『放之四海而皆准』的固定参数，必须根据数据物理尺寸动态调整！", 15, MED_BLUE, bold=True, sa=Pt(8))

# ================================================================
# SLIDE 7 — 训练策略
# ================================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, WHITE)
page_header(s7, "训练策略：组合损失函数", "Dice Loss + Cross Entropy 的黄金搭档")

# 总损失公式
box_f = add_rounded(s7, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2), RGBColor(0xE8, 0xF0, 0xF8), MED_BLUE, Pt(1))
tb_f = add_tb(s7, Inches(1.2), Inches(1.6), Inches(11.0), Inches(1.0))
tf = tb_f.text_frame
set_p(tf, "总损失  L = L_Dice + L_CE", 24, DARK_NAVY, bold=True, align=PP_ALIGN.CENTER, sa=Pt(4))
add_p(tf, "结合 Dice 对整体的『宏观把控』与 CE 对像素的『微观监督』", 16, MED_GRAY, align=PP_ALIGN.CENTER)

# 左：Dice Loss
box_d = add_rounded(s7, Inches(0.8), Inches(3.0), Inches(5.5), Inches(4.0), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GRN, Pt(1.5))
tb_d = add_tb(s7, Inches(1.0), Inches(3.2), Inches(5.0), Inches(3.6))
tf = tb_d.text_frame
set_p(tf, "Dice Loss", 22, ACCENT_GRN, bold=True, sa=Pt(8))
add_p(tf, "关注：预测区域与真实区域的『重合度』", 15, DARK_GRAY, sa=Pt(4))
add_p(tf, "本质：全局监督，整体度量分割质量", 15, DARK_GRAY, sa=Pt(4))
add_p(tf, "", 8, DARK_GRAY)
add_p(tf, "✓ 优势：天生抗不平衡", 15, ACCENT_GRN, bold=True, sa=Pt(4))
add_p(tf, "  目标区域极小（如血管）时仍有效", 14, DARK_GRAY, sa=Pt(4))
add_p(tf, "", 8, DARK_GRAY)
add_p(tf, "✗ 劣势：训练不稳定", 15, ACCENT_RED, bold=True, sa=Pt(4))
add_p(tf, "  梯度剧烈，易震荡", 14, DARK_GRAY, sa=Pt(4))

# 右：CE Loss
box_c = add_rounded(s7, Inches(7.0), Inches(3.0), Inches(5.5), Inches(4.0), RGBColor(0xFC, 0xE4, 0xE4), ACCENT_RED, Pt(1.5))
tb_c = add_tb(s7, Inches(7.2), Inches(3.2), Inches(5.0), Inches(3.6))
tf = tb_c.text_frame
set_p(tf, "Cross Entropy Loss", 22, ACCENT_RED, bold=True, sa=Pt(8))
add_p(tf, "关注：每个像素的分类『准确性』", 15, DARK_GRAY, sa=Pt(4))
add_p(tf, "本质：像素级监督，逐个分类判断", 15, DARK_GRAY, sa=Pt(4))
add_p(tf, "", 8, DARK_GRAY)
add_p(tf, "✓ 优势：训练稳定", 15, ACCENT_GRN, bold=True, sa=Pt(4))
add_p(tf, "  收敛平滑，梯度可控", 14, DARK_GRAY, sa=Pt(4))
add_p(tf, "", 8, DARK_GRAY)
add_p(tf, "✗ 劣势：偏向背景", 15, ACCENT_RED, bold=True, sa=Pt(4))
add_p(tf, "  前景过小时被背景主导", 14, DARK_GRAY, sa=Pt(4))

# ================================================================
# SLIDE 8 — 推理与后处理
# ================================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, WHITE)
page_header(s8, "推理与后处理：无缝拼接的艺术", "解决 Patch 预测的边缘不连续问题")

# 核心问题
box_q = add_rounded(s8, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0), RGBColor(0xFF, 0xF3, 0xE0), RGBColor(0xE6, 0x7E, 0x22), Pt(1))
tb_q = add_tb(s8, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.9))
tf = tb_q.text_frame
set_p(tf, "痛点：3D 图像切成 Patch 预测 → 边缘视野受限 → 拼图出现明显『断层疤痕』", 17, RGBColor(0xC6, 0x28, 0x28), bold=True)

# 解决方案
sol_items = [
    ("50% 重叠滑窗", ACCENT_TEAL, [
        "前一个块的『最边缘』= 后一个块的『正中心』",
        "中心视野好，预测精度极高",
        "例：第 127 号像素在第一块是边缘（预测差），",
        "  在第二块（64-191）是正中心（预测极准）",
    ]),
    ("高斯加权融合", MED_BLUE, [
        "Patch 中心 → 高权重（自信度 ~0.99）",
        "Patch 边缘 → 低权重（自信度 ~0.05）",
        "用最自信的中心区域覆盖边缘误差",
    ]),
    ("连通域后处理", RGBColor(0x7B, 0x2D, 0x8B), [
        "自动学习训练集标注的连通域数量",
        "如血管在训练集只有 1 个连通域",
        "测试时：保留最大连通域，孤立噪点→抹除",
    ]),
]
for i, (title, color, descs) in enumerate(sol_items):
    l = Inches(0.8 + i * 4.1)
    t = Inches(2.8)
    box = add_rounded(s8, l, t, Inches(3.8), Inches(4.2), RGBColor(0xF8, 0xF8, 0xF8), color, Pt(1.5))
    tb_t = add_tb(s8, l + Inches(0.15), t + Inches(0.15), Inches(3.5), Inches(0.45))
    set_p(tb_t.text_frame, f"方案 {i+1}：{title}", 18, color, bold=True)
    tb_d = add_tb(s8, l + Inches(0.15), t + Inches(0.7), Inches(3.5), Inches(3.3))
    tf = tb_d.text_frame
    for d in descs:
        add_p(tf, "• " + d, 14, DARK_GRAY, sa=Pt(4))

# ================================================================
# SLIDE 9 — Cascade 架构
# ================================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, WHITE)
page_header(s9, "3D U-Net Cascade：级联架构详解", "解决大图像爆显存 + 感受野不足的双重困境")

# Stage 1
box_s1 = add_rounded(s9, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), RGBColor(0xE8, 0xF0, 0xF8), MED_BLUE, Pt(1.5))
tb_s1 = add_tb(s9, Inches(1.0), Inches(1.6), Inches(5.0), Inches(2.3))
tf = tb_s1.text_frame
set_p(tf, "Stage 1：低分辨率全局定位", 20, DARK_NAVY, bold=True, sa=Pt(8))
for line in [
    "1. 原始高清大图 → 强行降采样（压扁缩小）",
    "2. 喂给第一个 3D U-Net，看清全局位置",
    "3. 输出『粗略的全局导航地图』(low res. seg.)",
    "4. 将导航地图放大回原始尺寸备用",
]:
    add_p(tf, line, 15, DARK_GRAY, sa=Pt(4))

# 箭头
arr = add_tb(s9, Inches(6.0), Inches(2.0), Inches(1.0), Inches(0.5))
set_p(arr.text_frame, "⬇ 融合 ⬇", 20, MED_BLUE, bold=True, align=PP_ALIGN.CENTER)

# Stage 2
box_s2 = add_rounded(s9, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5), RGBColor(0xFD, 0xF0, 0xE8), RGBColor(0xE6, 0x7E, 0x22), Pt(1.5))
tb_s2 = add_tb(s9, Inches(7.0), Inches(1.6), Inches(5.0), Inches(2.3))
tf = tb_s2.text_frame
set_p(tf, "Stage 2：高分辨率局部精细分割", 20, DARK_NAVY, bold=True, sa=Pt(8))
for line in [
    "1. 从原始高清大图切局部豆腐块（Patch）",
    "2. 将高清 Patch + Stage1 对应位置导航图 拼接",
    "3. 双重输入 → 第二个 3D U-Net",
    "4. 输出：精密的局部细节分割 (full res. seg.)",
]:
    add_p(tf, line, 15, DARK_GRAY, sa=Pt(4))

# 类比
box_ana = add_rounded(s9, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.5), RGBColor(0xF5, 0xF5, 0xF5), ACCENT_TEAL, Pt(1))
tb_ana = add_tb(s9, Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.2))
tf = tb_ana.text_frame
set_p(tf, "直观类比", 20, DARK_NAVY, bold=True, sa=Pt(8))
for line in [
    "Stage 1：先站到远处看全景，确定『目标大概在哪个区域』—— 全局视角，低分辨率。",
    "",
    "Stage 2：再拿着高清放大镜走近看细节，同时手里攥着 Stage 1 给的『全局 GPS 导航』",
    "         —— 局部高清 + 全局上下文，精度拉满。",
    "",
    "关键：两个网络共享相同的 3D U-Net 配置，区别仅在于输入分辨率和感受野。",
]:
    add_p(tf, line, 15, DARK_GRAY, sa=Pt(3))

# ================================================================
# SLIDE 10 — 效果与总结
# ================================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, WHITE)
page_header(s10, "实验结果与总结启示", "MSD 十项全能挑战赛 —— 最基础的 U-Net 击败所有魔改网络")

# MSD 结果框
box_msd = add_rounded(s10, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5), RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GRN, Pt(1.5))
tb_msd = add_tb(s10, Inches(1.0), Inches(1.6), Inches(11.3), Inches(2.3))
tf = tb_msd.text_frame
set_p(tf, "MSD 挑战赛结果", 22, ACCENT_GRN, bold=True, sa=Pt(8))
add_p(tf, "在医学图像分割十项全能（Medical Segmentation Decathlon）中，nnU-Net 以全自动配置的", 16, DARK_GRAY, sa=Pt(3))
add_p(tf, "最基础 U-Net，打败了所有当时最先进的『魔改』网络（Attention U-Net、Res-UNet 等）。", 16, DARK_GRAY, sa=Pt(3))
add_p(tf, "", 6, DARK_GRAY)
add_p(tf, "关键结论：在医学图像分割领域，数据适配流程 >> 网络结构创新", 17, MED_BLUE, bold=True, sa=Pt(4))

# 总结启示
box_sm = add_rounded(s10, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.6), RGBColor(0xFD, 0xF0, 0xE8), RGBColor(0xE6, 0x7E, 0x22), Pt(1))
tb_sm = add_tb(s10, Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.3))
tf = tb_sm.text_frame
set_p(tf, "总结与启示", 22, RGBColor(0xE6, 0x7E, 0x22), bold=True, sa=Pt(8))
points = [
    "❶ 接手新项目效果不好时 → 不要第一时间怀疑网络不够先进",
    "❷ 不要急着套用 Transformer / Attention → 先检查底层工程",
    "❸ 数据重采样做对了吗？归一化合理吗？Patch 切得合适吗？",
    "❹ 先把底层工程做到极致，建立坚如磐石的 Baseline",
    "❺ 这对深度学习效果有决定性的意义 —— 工程 > 魔改",
]
for p in points:
    add_p(tf, p, 16, DARK_GRAY, sa=Pt(6))

# ================================================================
# SAVE
# ================================================================
out = "d:/PythonProject/hello_agent/fastapi_/从U-Net到nnU-Net_组会汇报.pptx"
prs.save(out)
print(f"PPT saved: {out}")
